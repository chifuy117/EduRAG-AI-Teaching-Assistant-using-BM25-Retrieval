# ETL - Extract Module (Multi-Format Smart Loader)
# Handles extraction from: PDF, PPTX, DOCX, TXT, MD
# Automatically filters junk files, videos, images, compiled binaries
# Supports deep recursive directory traversal with metadata enrichment

import os
import re
import json
from datetime import datetime
from pathlib import Path


# ─────────────────────────────────────────────────────────────────────────────
# Document Model
# ─────────────────────────────────────────────────────────────────────────────

class Document:
    """Universal document class for all file types"""
    def __init__(self, page_content, metadata=None):
        self.page_content = page_content
        self.metadata = metadata or {}

    def __repr__(self):
        src = self.metadata.get('source_file', 'unknown')
        length = len(self.page_content)
        return f"<Document source='{src}' chars={length}>"


# ─────────────────────────────────────────────────────────────────────────────
# File Classification Constants
# ─────────────────────────────────────────────────────────────────────────────

# Extensions we CAN process
PROCESSABLE_EXTENSIONS = {
    '.pdf',     # PDF documents
    '.pptx',    # PowerPoint presentations
    '.docx',    # Word documents
    '.txt',     # Plain text files
    '.md',      # Markdown files
}

# Extensions we ALWAYS skip (junk / binary / media)
SKIP_EXTENSIONS = {
    # Videos
    '.mp4', '.avi', '.mkv', '.mov', '.wmv', '.flv', '.webm', '.m4v',
    # Images (unless OCR is enabled)
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.svg', '.ico', '.webp', '.tiff',
    # Compiled / binary
    '.exe', '.o', '.obj', '.dll', '.so', '.class', '.pyc', '.pyo',
    # Archives
    '.zip', '.rar', '.7z', '.tar', '.gz', '.bz2',
    # Source code (not useful for RAG text)
    '.c', '.cpp', '.h', '.java', '.py', '.js', '.ts', '.html', '.css',
    # Database / system
    '.db', '.sqlite', '.log', '.bak', '.tmp',
}

# Filename patterns that indicate junk / temp files
JUNK_PATTERNS = [
    r'^~\$',               # Office temp files (~$filename)
    r'^\.~',               # Hidden temp files
    r'\.tmp$',             # Temp files
    r'^Thumbs\.db$',       # Windows thumbnail cache
    r'^\.DS_Store$',       # macOS metadata
    r'^desktop\.ini$',     # Windows folder config
    r'^__pycache__$',      # Python cache
]


# ─────────────────────────────────────────────────────────────────────────────
# Individual File Extractors
# ─────────────────────────────────────────────────────────────────────────────

def extract_pdf(file_path):
    """Extract text from PDF files"""
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        pages_text = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages_text.append(text.strip())

        if not pages_text:
            return []

        full_text = "\n\n".join(pages_text)
        doc = Document(
            page_content=full_text,
            metadata={
                'source': file_path,
                'source_type': 'pdf',
                'pages': len(reader.pages),
                'pages_with_text': len(pages_text),
            }
        )
        return [doc]
    except Exception as e:
        print(f"    ⚠️  Error reading PDF {os.path.basename(file_path)}: {e}")
        return []


def extract_pptx(file_path):
    """Extract text from PowerPoint files"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_content.append(text)

                # Also extract text from tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_text.append(cell_text)
                        if row_text:
                            slide_content.append(" | ".join(row_text))

            if slide_content:
                slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_content))

        if not slides_text:
            return []

        full_text = "\n\n".join(slides_text)
        doc = Document(
            page_content=full_text,
            metadata={
                'source': file_path,
                'source_type': 'pptx',
                'total_slides': len(prs.slides),
                'slides_with_text': len(slides_text),
            }
        )
        return [doc]
    except Exception as e:
        print(f"    ⚠️  Error reading PPTX {os.path.basename(file_path)}: {e}")
        return []


def extract_docx(file_path):
    """Extract text from Word documents"""
    try:
        from docx import Document as DocxDocument
        doc_file = DocxDocument(file_path)
        paragraphs = []

        for para in doc_file.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        # Also extract from tables
        for table in doc_file.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                if row_text:
                    paragraphs.append(" | ".join(row_text))

        if not paragraphs:
            return []

        full_text = "\n\n".join(paragraphs)
        doc = Document(
            page_content=full_text,
            metadata={
                'source': file_path,
                'source_type': 'docx',
                'paragraphs': len(paragraphs),
            }
        )
        return [doc]
    except Exception as e:
        print(f"    ⚠️  Error reading DOCX {os.path.basename(file_path)}: {e}")
        return []


def extract_txt(file_path):
    """Extract text from plain text / markdown files"""
    try:
        # Try UTF-8 first, fall back to latin-1
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read().strip()
                break
            except UnicodeDecodeError:
                continue
        else:
            return []

        if not text or len(text) < 10:
            return []

        doc = Document(
            page_content=text,
            metadata={
                'source': file_path,
                'source_type': 'txt',
                'char_count': len(text),
                'line_count': text.count('\n') + 1,
            }
        )
        return [doc]
    except Exception as e:
        print(f"    ⚠️  Error reading TXT {os.path.basename(file_path)}: {e}")
        return []


# ─────────────────────────────────────────────────────────────────────────────
# File Classification Helpers
# ─────────────────────────────────────────────────────────────────────────────

def is_junk_file(filename):
    """Check if a file is a junk/temp file that should be skipped"""
    for pattern in JUNK_PATTERNS:
        if re.match(pattern, filename, re.IGNORECASE):
            return True
    return False


def classify_file(file_path):
    """
    Classify a file as 'process', 'skip', or 'junk'
    Returns (classification, extension)
    """
    filename = os.path.basename(file_path)
    ext = os.path.splitext(filename)[1].lower()

    if is_junk_file(filename):
        return 'junk', ext

    if ext in PROCESSABLE_EXTENSIONS:
        return 'process', ext

    if ext in SKIP_EXTENSIONS:
        return 'skip', ext

    # Unknown extension — skip by default
    return 'skip', ext


def detect_subject(file_path, base_path):
    """
    Auto-detect subject from path structure.
    Looks for known subject folder names in the path.
    """
    # Get relative path from base
    try:
        rel_path = os.path.relpath(file_path, base_path)
    except ValueError:
        rel_path = file_path

    from pathlib import Path
    parts = Path(rel_path).parts

    # The subject is the immediate child of the data directory
    if len(parts) > 1:
        return parts[0] # Return exactly the folder name (e.g., 'MachineLearning')

    return 'General'


def detect_unit(file_path):
    """Auto-detect unit number from path or filename"""
    full_path = file_path.replace('\\', '/')

    # Look for unit patterns in path and filename
    patterns = [
        r'[Uu]nit[_\s-]*(\d+)',
        r'UNIT[_\s-]*(\d+)',
        r'[Uu](\d+)',
    ]

    for pattern in patterns:
        match = re.search(pattern, full_path)
        if match:
            return f"Unit_{match.group(1)}"

    return 'Unknown'


# ─────────────────────────────────────────────────────────────────────────────
# Extraction Router
# ─────────────────────────────────────────────────────────────────────────────

# Map extensions to their extractor functions
EXTRACTOR_MAP = {
    '.pdf':  extract_pdf,
    '.pptx': extract_pptx,
    '.docx': extract_docx,
    '.txt':  extract_txt,
    '.md':   extract_txt,       # Markdown uses same loader as txt
}


def extract_file(file_path):
    """Route a file to the correct extractor based on extension"""
    ext = os.path.splitext(file_path)[1].lower()
    extractor = EXTRACTOR_MAP.get(ext)

    if extractor is None:
        return []

    return extractor(file_path)


# ─────────────────────────────────────────────────────────────────────────────
# Main DataExtractor Class
# ─────────────────────────────────────────────────────────────────────────────

class DataExtractor:
    """
    Smart multi-format data extractor with:
    - Recursive directory traversal
    - Automatic junk file filtering
    - Multi-format support (PDF, PPTX, DOCX, TXT, MD)
    - Rich metadata enrichment (subject, unit, source info)
    - Processing statistics and reporting
    """

    def __init__(self, data_dir="data"):
        self.data_dir = data_dir
        os.makedirs(data_dir, exist_ok=True)

        # Processing statistics
        self.stats = {
            'total_files_scanned': 0,
            'files_processed': 0,
            'files_skipped': 0,
            'junk_removed': 0,
            'documents_created': 0,
            'by_type': {},
            'by_subject': {},
            'errors': [],
            'skipped_files': [],
            'junk_files': [],
        }

    def reset_stats(self):
        """Reset processing statistics"""
        self.stats = {
            'total_files_scanned': 0,
            'files_processed': 0,
            'files_skipped': 0,
            'junk_removed': 0,
            'documents_created': 0,
            'by_type': {},
            'by_subject': {},
            'errors': [],
            'skipped_files': [],
            'junk_files': [],
        }

    def extract_from_pdf(self, pdf_path):
        """Extract text from a single PDF file (backward compatibility)"""
        return extract_pdf(pdf_path)

    def extract_single_file(self, file_path, base_path=None):
        """
        Extract content from a single file with full metadata enrichment
        """
        if base_path is None:
            base_path = os.path.dirname(file_path)

        filename = os.path.basename(file_path)
        classification, ext = classify_file(file_path)

        self.stats['total_files_scanned'] += 1

        # Handle junk files
        if classification == 'junk':
            self.stats['junk_removed'] += 1
            self.stats['junk_files'].append(filename)
            return []

        # Handle skipped files
        if classification == 'skip':
            self.stats['files_skipped'] += 1
            self.stats['skipped_files'].append(filename)
            return []

        # Process the file
        docs = extract_file(file_path)

        if not docs:
            self.stats['files_skipped'] += 1
            return []

        # Detect metadata
        subject = detect_subject(file_path, base_path)
        unit = detect_unit(file_path)

        # Enrich each document with metadata
        for doc in docs:
            doc.metadata.update({
                'source_file': filename,
                'subject': subject,
                'unit': unit,
                'file_extension': ext,
                'file_size_bytes': os.path.getsize(file_path),
                'extraction_timestamp': datetime.now().isoformat(),
                'doc_id': f"{subject}_{unit}_{filename}_{hash(doc.page_content[:100])}",
            })

        # Update stats
        self.stats['files_processed'] += 1
        self.stats['documents_created'] += len(docs)

        ext_key = ext.lstrip('.')
        self.stats['by_type'][ext_key] = self.stats['by_type'].get(ext_key, 0) + 1
        self.stats['by_subject'][subject] = self.stats['by_subject'].get(subject, 0) + len(docs)

        return docs

    def extract_from_directory(self, dir_path, recursive=True):
        """
        Extract all processable files from a directory.
        Supports both flat and deeply nested folder structures.
        """
        if not os.path.exists(dir_path):
            print(f"❌ Directory not found: {dir_path}")
            return []

        all_documents = []
        base_path = dir_path

        print(f"\n{'='*60}")
        print(f"📂 Scanning: {os.path.abspath(dir_path)}")
        print(f"{'='*60}")

        if recursive:
            for root, dirs, files in os.walk(dir_path):
                # Skip hidden directories and __pycache__
                dirs[:] = [d for d in dirs if not d.startswith('.') and d != '__pycache__' and d != 'vector_db']

                # Get relative path for display
                rel_root = os.path.relpath(root, dir_path)
                if rel_root == '.':
                    rel_root = '(root)'

                processable = [f for f in files if classify_file(os.path.join(root, f))[0] == 'process']

                if processable:
                    print(f"\n📁 {rel_root}/ ({len(processable)} processable files)")

                for filename in sorted(files):
                    file_path = os.path.join(root, filename)
                    classification, ext = classify_file(file_path)

                    if classification == 'process':
                        print(f"    ✅ {filename} [{ext}]")
                        docs = self.extract_single_file(file_path, base_path)
                        all_documents.extend(docs)
                    elif classification == 'junk':
                        print(f"    🗑️  {filename} [JUNK - removed]")
                        self.stats['junk_removed'] += 1
                        self.stats['junk_files'].append(filename)
                        self.stats['total_files_scanned'] += 1
                    else:
                        # Silently skip for cleaner output
                        self.stats['files_skipped'] += 1
                        self.stats['skipped_files'].append(filename)
                        self.stats['total_files_scanned'] += 1
        else:
            # Flat scan (non-recursive)
            for filename in sorted(os.listdir(dir_path)):
                file_path = os.path.join(dir_path, filename)
                if os.path.isfile(file_path):
                    docs = self.extract_single_file(file_path, base_path)
                    all_documents.extend(docs)

        return all_documents

    def extract_all_sources(self, pdf_dir=None, data_dir=None):
        """
        Extract from all configured sources (main entry point).
        Accepts either pdf_dir (backward compat) or data_dir.
        """
        target_dir = data_dir or pdf_dir or self.data_dir
        all_documents = []

        if target_dir and os.path.exists(target_dir):
            docs = self.extract_from_directory(target_dir, recursive=True)
            all_documents.extend(docs)
        else:
            print(f"❌ Directory not found: {target_dir}")

        return all_documents

    def print_stats(self):
        """Print a comprehensive processing report"""
        s = self.stats

        print(f"\n{'='*60}")
        print(f"📊 EXTRACTION REPORT")
        print(f"{'='*60}")
        print(f"  Total files scanned  : {s['total_files_scanned']}")
        print(f"  Files processed      : {s['files_processed']}")
        print(f"  Files skipped        : {s['files_skipped']}")
        print(f"  Junk files removed   : {s['junk_removed']}")
        print(f"  Documents created    : {s['documents_created']}")

        if s['by_type']:
            print(f"\n  📄 By File Type:")
            for ftype, count in sorted(s['by_type'].items()):
                print(f"      .{ftype:6s} → {count} file(s)")

        if s['by_subject']:
            print(f"\n  🎓 By Subject:")
            for subject, count in sorted(s['by_subject'].items()):
                print(f"      {subject:20s} → {count} document(s)")

        if s['junk_files']:
            print(f"\n  🗑️  Junk Files Filtered:")
            for jf in s['junk_files']:
                print(f"      - {jf}")

        if s['errors']:
            print(f"\n  ❌ Errors:")
            for err in s['errors']:
                print(f"      - {err}")

        print(f"{'='*60}\n")

    def get_stats(self):
        """Return processing stats as a dictionary"""
        return self.stats.copy()


# ─────────────────────────────────────────────────────────────────────────────
# Standalone Usage
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import sys

    print("🔧 Multi-Format Data Extractor")
    print("=" * 40)

    if len(sys.argv) > 1:
        target = sys.argv[1]
    else:
        target = "data"

    extractor = DataExtractor()
    documents = extractor.extract_from_directory(target)
    extractor.print_stats()

    if documents:
        print(f"\n📄 Sample document preview:")
        print(f"   Source: {documents[0].metadata.get('source_file', 'N/A')}")
        print(f"   Subject: {documents[0].metadata.get('subject', 'N/A')}")
        print(f"   Unit: {documents[0].metadata.get('unit', 'N/A')}")
        print(f"   Content preview: {documents[0].page_content[:200]}...")